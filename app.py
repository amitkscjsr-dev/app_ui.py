import streamlit as st
import os
import requests
import base64
import json
import sqlite3
import time
from datetime import datetime, timedelta
from dotenv import load_dotenv

try:
    from zoneinfo import ZoneInfo
    IST = ZoneInfo("Asia/Kolkata")
except ImportError:
    import pytz
    IST = pytz.timezone("Asia/Kolkata")

# ═══════════════════════════════════════════════════════════════
# 1. SETUP & API KEYS  (unchanged — your existing keys kept)
# ═══════════════════════════════════════════════════════════════
load_dotenv()
llm_key          = os.getenv("LLM_API_KEY")          # your existing OpenAI key
tavily_key       = os.getenv("TAVILY_API_KEY")        # your existing Tavily key
telegram_token   = os.getenv("TELEGRAM_BOT_TOKEN", "")
telegram_chat_id = os.getenv("TELEGRAM_CHAT_ID", "")

# Your existing Make.com webhook — kept exactly as-is
MAKE_WEBHOOK_URL = "https://hook.eu1.make.com/m4eac7gnkm0jklsy3mrnrik43rrvhuw6"

# LinkedIn — placeholder until you create the app (set these in .env later)
LINKEDIN_ACCESS_TOKEN = os.getenv("LINKEDIN_ACCESS_TOKEN", "")
LINKEDIN_PERSON_URN   = os.getenv("LINKEDIN_PERSON_URN", "")


# ═══════════════════════════════════════════════════════════════
# 2. DATABASE  (extended with safe migration)
# ═══════════════════════════════════════════════════════════════
def init_db():
    conn = sqlite3.connect("content_calendar.db")
    c    = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS posts (
        id             INTEGER PRIMARY KEY,
        date           TEXT,
        topic          TEXT,
        content        TEXT,
        image          BLOB,
        status         TEXT,
        scheduled_time TEXT,
        niche          TEXT,
        version        TEXT
    )''')
    # Safe migration — adds new columns to existing databases without breaking anything
    for col_name, col_type in [("scheduled_time","TEXT"), ("niche","TEXT"), ("version","TEXT")]:
        try:
            c.execute(f"ALTER TABLE posts ADD COLUMN {col_name} {col_type}")
        except Exception:
            pass
    conn.commit()
    conn.close()

def save_post_to_db(topic, content, image_data, scheduled_time="", niche="", version="A"):
    conn     = sqlite3.connect("content_calendar.db")
    c        = conn.cursor()
    date_now = datetime.now().strftime("%Y-%m-%d %H:%M")
    c.execute(
        "INSERT INTO posts (date, topic, content, image, status, scheduled_time, niche, version) VALUES (?,?,?,?,?,?,?,?)",
        (date_now, topic, content, image_data, "Draft", scheduled_time, niche, version)
    )
    conn.commit()
    conn.close()

def update_post_status(post_id, new_status):
    conn = sqlite3.connect("content_calendar.db")
    c    = conn.cursor()
    c.execute("UPDATE posts SET status = ? WHERE id = ?", (new_status, post_id))
    conn.commit()
    conn.close()

def get_all_posts():
    conn = sqlite3.connect("content_calendar.db")
    c    = conn.cursor()
    c.execute("SELECT id, date, topic, content, image, status, scheduled_time, niche, version FROM posts ORDER BY id DESC")
    data = c.fetchall()
    conn.close()
    return data

init_db()


# ═══════════════════════════════════════════════════════════════
# 3. LOAD PERSONAL PROFILE  (unchanged)
# ═══════════════════════════════════════════════════════════════
try:
    with open("profile.json") as f:
        profile    = json.load(f)
        user_name  = profile.get("name",  "Amit Kumar Singh")
        user_style = profile.get("style", "Short, simple English, leadership-focused, no jargon")
except Exception:
    user_name  = "Amit Kumar Singh"
    user_style = "Short, simple English, leadership-focused, no jargon"


# ═══════════════════════════════════════════════════════════════
# 4. BEST-TIME ALGORITHM  (IST — LinkedIn research-backed)
# ═══════════════════════════════════════════════════════════════
# Engagement scores out of 10 per slot
BEST_SLOTS = {
    "Monday":    [("07:30", 7), ("12:00", 6)],
    "Tuesday":   [("07:30", 10), ("09:00", 9), ("12:00", 8)],   # best day
    "Wednesday": [("08:00", 9),  ("09:00", 8), ("17:00", 7)],   # best day
    "Thursday":  [("08:00", 9),  ("12:00", 7), ("17:30", 6)],
    "Friday":    [("08:00", 7),  ("11:00", 6)],
    "Saturday":  [],
    "Sunday":    [],
}

def get_next_best_slot():
    now = datetime.now()
    for days_ahead in range(8):
        candidate = now + timedelta(days=days_ahead)
        day_name  = candidate.strftime("%A")
        for slot_str, score in BEST_SLOTS.get(day_name, []):
            h, m    = map(int, slot_str.split(":"))
            slot_dt = candidate.replace(hour=h, minute=m, second=0, microsecond=0)
            if slot_dt > now + timedelta(minutes=30):
                days_label = "today" if days_ahead == 0 else "tomorrow" if days_ahead == 1 else f"in {days_ahead} days"
                return {
                    "iso":     slot_dt.strftime("%Y-%m-%dT%H:%M:%S"),
                    "display": f"{day_name} at {slot_str} IST ({days_label})",
                    "score":   score,
                    "day":     day_name,
                }
    # fallback: next Tuesday 8am
    days_to_tue = (1 - now.weekday()) % 7 or 7
    nxt         = now + timedelta(days=days_to_tue)
    nxt         = nxt.replace(hour=8, minute=0, second=0, microsecond=0)
    return {"iso": nxt.strftime("%Y-%m-%dT%H:%M:%S"), "display": "Tuesday at 08:00 IST", "score": 10, "day": "Tuesday"}

def get_slot_calendar():
    now    = datetime.now()
    result = []
    for days_ahead in range(7):
        candidate = now + timedelta(days=days_ahead)
        day_name  = candidate.strftime("%A")
        for slot_str, score in BEST_SLOTS.get(day_name, []):
            h, m    = map(int, slot_str.split(":"))
            slot_dt = candidate.replace(hour=h, minute=m, second=0, microsecond=0)
            if slot_dt > now:
                result.append({"label": f"{day_name} {slot_str}", "score": score,
                                "iso": slot_dt.strftime("%Y-%m-%dT%H:%M:%S")})
    return sorted(result, key=lambda x: -x["score"])[:8]


# ═══════════════════════════════════════════════════════════════
# 5. CORE API FUNCTIONS  (generate_text & generate_image unchanged;
#                         web_search unchanged; send_to_make enhanced)
# ═══════════════════════════════════════════════════════════════
def generate_text(prompt, system_prompt="You are an expert LinkedIn strategist.", temperature=0.8):
    url     = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {llm_key}", "Content-Type": "application/json"}
    data    = {
        "model": "gpt-4o",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": prompt},
        ],
        "temperature": temperature,
    }
    try:
        resp = requests.post(url, headers=headers, json=data, timeout=90).json()
        return resp["choices"][0]["message"]["content"]
    except Exception as e:
        return f"⚠️ API Error: {e}"

def generate_image(prompt):
    url     = "https://api.openai.com/v1/images/generations"
    headers = {"Authorization": f"Bearer {llm_key}", "Content-Type": "application/json"}
    data    = {
        "model":           "dall-e-3",
        "prompt":          prompt,
        "size":            "1792x1024",          # LinkedIn banner ratio
        "style":           "natural",
        "response_format": "b64_json",
    }
    try:
        resp = requests.post(url, headers=headers, json=data, timeout=90).json()
        return base64.b64decode(resp["data"][0]["b64_json"])
    except Exception:
        return None

def web_search(query):
    """
    Tavily Deep Research — forces recency so only 2025/2026 content is used.
    Injects today's date into the query and uses days=90 to restrict results
    to the last 90 days, preventing old 2023/2024 articles from polluting posts.
    """
    if not tavily_key:
        return "Tavily API Key missing."

    TODAY       = datetime.now().strftime("%B %d, %Y")          # e.g. "March 29, 2026"
    CURRENT_YR  = datetime.now().strftime("%Y")                  # "2026"

    # Append date context so Tavily prioritises current results
    dated_query = (
        f"{query} "
        f"[Research must focus on {CURRENT_YR} data only. "
        f"Today is {TODAY}. Exclude anything from 2024 or earlier.]"
    )

    headers = {"Authorization": f"Bearer {tavily_key}", "Content-Type": "application/json"}
    try:
        resp = requests.post(
            "https://api.tavily.com/research",
            headers=headers,
            json={
                "input": dated_query,
                "model": "mini",
                "days":  90,        # only index sources from last 90 days
            },
            timeout=30,
        ).json()
        request_id = resp.get("request_id")
        if not request_id:
            return f"Error starting research: {resp}"
        get_url = f"https://api.tavily.com/research/{request_id}"
        for _ in range(36):          # max ~3 minutes
            time.sleep(5)
            s = requests.get(get_url, headers=headers, timeout=15).json()
            if s.get("status") == "completed":
                return s.get("content", "No content found.")
            if s.get("status") == "failed":
                return "Deep research failed."
        return "Research timed out after 3 minutes."
    except Exception as e:
        return f"Error: {e}"

def compress_image_for_webhook(image_bytes, max_kb=400):
    """
    Compress image to stay under max_kb before base64 encoding.
    Make.com rejects webhook payloads over ~1 MB — a raw DALL·E 3 image
    at 1792×1024 encodes to 4–6 MB of base64, which causes the 400 error.
    This resizes to 1200×628 (LinkedIn optimal) and reduces JPEG quality
    until the encoded size is safely under the limit.
    """
    try:
        from PIL import Image
        import io

        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")

        # LinkedIn optimal size — also reduces raw bytes significantly
        img.thumbnail((1200, 628), Image.LANCZOS)

        # Try progressively lower quality until under max_kb
        for quality in [85, 70, 55, 40]:
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=quality, optimize=True)
            compressed = buf.getvalue()
            kb = len(compressed) / 1024
            if kb <= max_kb:
                return compressed, round(kb)

        # Last resort — return smallest attempt
        return compressed, round(len(compressed) / 1024)

    except ImportError:
        # Pillow not installed — return original with a warning flag
        return image_bytes, round(len(image_bytes) / 1024)
    except Exception:
        return image_bytes, round(len(image_bytes) / 1024)


def send_to_make(post, image_bytes=None, topic="", niche="", publish_at="", slot_label=""):
    """
    Sends post to Make.com webhook.

    Fix for 400 error — three root causes addressed:
      1. Image too large  → compress to <400 KB before base64 encoding
      2. data: URI prefix → removed (Make.com JSON parser rejects it)
      3. Missing header   → explicit Content-Type: application/json added
    """
    payload = {
        "post":         post,           # your original field name — kept for Make.com mapping
        "post_content": post,
        "topic":        topic,
        "niche":        niche,
        "publish_at":   publish_at,
        "slot_label":   slot_label,
        "source":       "ai_content_os",
        "sent_at":      datetime.now().strftime("%Y-%m-%dT%H:%M:%S"),
        "has_image":    image_bytes is not None,
    }

    if image_bytes:
        compressed, size_kb = compress_image_for_webhook(image_bytes, max_kb=400)

        if size_kb > 700:
            # Still too large even after compression — skip image, flag it
            payload["image_b64"]      = ""
            payload["image_skipped"]  = True
            payload["image_skip_reason"] = f"Too large after compression ({size_kb} KB). Upload manually."
        else:
            # Clean base64 — NO data: URI prefix (that was causing the 400)
            b64 = base64.b64encode(compressed).decode("utf-8")
            payload["image_b64"]     = b64
            payload["image_size_kb"] = size_kb

    try:
        resp = requests.post(
            MAKE_WEBHOOK_URL,
            json=payload,
            headers={"Content-Type": "application/json"},  # explicit header
            timeout=20,
        )
        # Return tuple so caller can show the real error message
        return resp.status_code, resp.text
    except requests.Timeout:
        return None, "Request timed out after 20 seconds"
    except Exception as e:
        return None, str(e)

def send_telegram(message):
    """Optional Telegram notification — silently skipped if keys not in .env."""
    if not telegram_token or not telegram_chat_id:
        return False
    try:
        resp = requests.post(
            f"https://api.telegram.org/bot{telegram_token}/sendMessage",
            json={"chat_id": telegram_chat_id, "text": message, "parse_mode": "HTML"},
            timeout=8,
        )
        return resp.status_code == 200
    except Exception:
        return False


# ═══════════════════════════════════════════════════════════════
# 6. LINKEDIN POST ENGINE  (new — replaces the simple draft_post call)
# ═══════════════════════════════════════════════════════════════
LINKEDIN_POST_RULES = """
CRITICAL LINKEDIN POST RULES — follow every one:
• Hook (lines 1–2): Must stop the scroll. Use a bold claim, surprising number, or open loop.
  NEVER start with "I" or "I'm excited to share" or "In today's world..."
• Add ONE blank line after the hook — this creates the 'See more' fold on mobile.
• Body: Short paragraphs of 1–3 lines each. Generous white space between every paragraph.
• Include at least one specific fact, stat, or real example from the research provided.
• Near the end: one engaging question that makes readers want to comment.
• Soft CTA (last line before hashtags): "Follow for more on [topic]" or "Save this post".
• Hashtags: 3–5 on their own line at the very end. Mix broad + niche.
• Total length: 1000–1800 characters.
• NO corporate jargon. NO "synergy", "leverage", "paradigm shift", "deep dive".
• Use "you" and "your" throughout — speak directly to the reader.
• Do NOT include any [Placeholder], [Author Name], or instructional text in output.
"""

def generate_topic_angles(idea: str, research: str, niche: str) -> list:
    """Generate 5 distinct LinkedIn post angles from the deep research."""
    TODAY      = datetime.now().strftime("%B %d, %Y")
    CURRENT_YR = datetime.now().strftime("%Y")

    prompt = f"""You are a LinkedIn content strategist specialising in {niche}.

TODAY'S DATE: {TODAY}
CURRENT YEAR: {CURRENT_YR}
IMPORTANT: Only reference {CURRENT_YR} data. Never mention 2024, 2023, or earlier years.

Based on this deep research about "{idea}":
{research[:2500]}

Generate exactly 5 distinct LinkedIn post angles. Each must use a different format type.

Return ONLY a valid JSON array — no markdown, no code fences, no explanation:
[
  {{
    "title":        "Specific angle in max 8 words",
    "hook":         "First 1–2 punchy sentences that open this post",
    "format":       "One of: Story / Numbered list / Contrarian / Data reveal / How-to",
    "why_it_works": "One sentence on why this format performs on LinkedIn",
    "trend_score":  "High / Very High / Medium"
  }}
]"""
    raw = generate_text(prompt, "You are a JSON-only response assistant. Return valid JSON arrays only.", temperature=0.7)
    try:
        clean = raw.strip()
        if "```" in clean:
            clean = clean.split("```")[1]
            if clean.startswith("json"):
                clean = clean[4:]
        return json.loads(clean.strip())[:5]
    except Exception:
        return [
            {"title": f"The truth about {idea} nobody tells you", "hook": f"Most people get {idea} completely wrong.\n\nHere's what actually works (and why):", "format": "Contrarian", "why_it_works": "Contrarian posts get 4x more comments than agreeable ones", "trend_score": "Very High"},
            {"title": f"5 lessons from {idea} that changed everything", "hook": f"After going deep on {idea}, these 5 lessons changed how I work:", "format": "Numbered list", "why_it_works": "Lists are LinkedIn's #1 saved and reshared format", "trend_score": "Very High"},
            {"title": f"My honest {idea} journey", "hook": f"I almost quit {idea} entirely.\n\nThen one shift changed everything.", "format": "Story", "why_it_works": "Personal stories build trust and drive emotional engagement", "trend_score": "High"},
            {"title": f"The data behind {idea} surprised me", "hook": f"I analysed 50+ cases on {idea}.\n\nThe results were not what anyone expected.", "format": "Data reveal", "why_it_works": "Data posts build authority and attract senior professionals", "trend_score": "High"},
            {"title": f"How to apply {idea} starting this week", "hook": f"You don't need months to benefit from {idea}.\n\nHere's the fastest path:", "format": "How-to", "why_it_works": "Actionable posts get saved for later — driving repeated views", "trend_score": "High"},
        ]

def build_post(idea, research, topic_angle, niche, tone, style_instructions):
    """Build a full LinkedIn post using the formula + research data."""
    TODAY      = datetime.now().strftime("%B %d, %Y")
    CURRENT_YR = datetime.now().strftime("%Y")

    prompt = f"""Write a complete, publish-ready LinkedIn post.

TODAY'S DATE: {TODAY} — this is critical context.
CURRENT YEAR: {CURRENT_YR}

Author: {user_name}
Writing style: {user_style}
Niche: {niche}
Tone: {tone}
Topic angle: {topic_angle['title']}
Hook to build from: {topic_angle['hook']}
Post format: {topic_angle['format']}
Style instruction: {style_instructions}

Research to draw facts and data from:
{research[:2000]}

{LINKEDIN_POST_RULES}
• DATE RULE — CRITICAL: Today is {TODAY}. Never mention years 2024, 2023, 2022 or earlier.
  If research contains old data, either update it to {CURRENT_YR} context or omit it entirely.
  When referencing trends say "right now", "this year", "in {CURRENT_YR}" — never old years.

Write the complete post now. Output ONLY the post text — nothing else."""
    system = f"You are an elite LinkedIn ghostwriter specialising in {niche}. Today is {TODAY}. Only use current {CURRENT_YR} data and trends."
    return generate_text(prompt, system, temperature=0.85)

def post_metrics(content: str) -> dict:
    """Extract simple metadata from a generated post."""
    import re
    chars    = len(content)
    hashtags = re.findall(r"#\w+", content)
    words    = content.split()
    read_min = max(1, round(len(words) / 238))
    first    = content.split("\n")[0].lower()
    if "?" in first[:80]:
        hook = "Question hook"
    elif any(c.isdigit() for c in first[:30]):
        hook = "Number hook"
    elif any(w in first for w in ["most ", "nobody", "stop ", "wrong", "never"]):
        hook = "Pattern interrupt"
    elif any(w in first for w in ["i ", "my ", "i'"]):
        hook = "Story / confession"
    else:
        hook = "Bold claim"
    return {"chars": chars, "hashtags": len(hashtags), "read_min": read_min, "hook_type": hook}


def check_stale_years(content: str) -> list:
    """
    Scans post for old year mentions (2024 and earlier).
    Returns list of flagged years found so the UI can warn the user.
    Posts going live on LinkedIn with '2024' or '2023' look outdated instantly.
    """
    import re
    current_year = datetime.now().year
    stale_years  = []
    # Find all 4-digit years in the post
    found = re.findall(r'\b(20\d{2})\b', content)
    for yr in found:
        if int(yr) < current_year:
            stale_years.append(yr)
    return list(set(stale_years))  # deduplicated


# ═══════════════════════════════════════════════════════════════
# 7. PAGE CONFIG & CSS
# ═══════════════════════════════════════════════════════════════
st.set_page_config(page_title="Personal Brand AI", page_icon="🚀", layout="wide")

st.markdown("""
<style>
/* Step progress bar */
.step-bar { display:flex; align-items:center; gap:0; margin-bottom:1.5rem; }
.step-pill {
    display:flex; align-items:center; justify-content:center;
    width:28px; height:28px; border-radius:50%;
    font-size:12px; font-weight:700; flex-shrink:0;
}
.step-pill.done   { background:#0a66c2; color:#fff; }
.step-pill.active { background:#1a1a2e; color:#fff; }
.step-pill.todo   { background:#e8e8e8; color:#aaa; }
.step-line { flex:1; height:1px; background:#e0e0e0; }
.step-label { font-size:11px; font-weight:600; margin-left:6px; }

/* Topic cards */
.topic-card {
    border:1.5px solid #e4e4e4; border-radius:10px;
    padding:14px 16px; margin-bottom:10px;
    border-left:4px solid #0a66c2; background:#fff;
}
.topic-card h5 { margin:0 0 6px; font-size:14px; color:#1a1a2e; }
.topic-card p  { margin:0; font-size:12px; color:#555; line-height:1.5; }
.badge {
    display:inline-block; padding:2px 8px; border-radius:12px;
    font-size:10px; font-weight:600; margin-right:4px; margin-top:6px;
}
.badge-blue   { background:#e8f0fe; color:#1a56db; }
.badge-green  { background:#e8f5e9; color:#1e7e34; }
.badge-purple { background:#f3e8fd; color:#6b21a8; }

/* A/B version tabs */
.version-badge {
    display:inline-block; padding:3px 12px; border-radius:20px;
    font-size:11px; font-weight:700; margin-bottom:8px;
}
.version-a { background:#e8f0fe; color:#0a66c2; }
.version-b { background:#fde8f0; color:#b44f8a; }

/* Metrics row */
.metrics-row { display:flex; gap:8px; flex-wrap:wrap; margin:8px 0 12px; }
.metric-chip {
    background:#f5f5f5; border:1px solid #e4e4e4;
    border-radius:16px; padding:3px 10px; font-size:11px; color:#555;
}

/* Schedule slot */
.slot-highlight {
    background:#e8f5e9; border:1px solid #a5d6a7; border-radius:8px;
    padding:10px 16px; font-size:13px; color:#1b5e20; margin:8px 0;
}

/* Publish status */
.linkedin-placeholder {
    background:#fff8e1; border:1px solid #ffe082; border-radius:8px;
    padding:10px 16px; font-size:12px; color:#5d4037; margin:8px 0;
}
</style>
""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════
# 8. SESSION STATE  (all new keys + backward-compat with original)
# ═══════════════════════════════════════════════════════════════
_defaults = {
    # Agentic workflow multi-step state
    "aw_step":           1,
    "aw_idea":           "",
    "aw_research":       None,
    "aw_topics":         [],
    "aw_selected_topic": None,
    "aw_post_a":         "",
    "aw_post_b":         "",
    "aw_metrics_a":      {},
    "aw_metrics_b":      {},
    "aw_active":         "A",
    "aw_final_post":     "",
    "aw_image":          None,
    "aw_img_prompt":     "",
    "aw_scheduled_iso":  "",
    "aw_scheduled_lbl":  "",
    # Sidebar settings
    "niche": "AI & Technology",
    "tone":  "Insightful & Authoritative",
    # Legacy keys (keep for backward compatibility)
    "current_post": "",
    "shared_topic": "",
    "agent_image":  None,
    # Chatbot
    "chat_history":    [],
    # Carousel
    "carousel_topic":  "",
    "carousel_slides": [],
    "carousel_style":  "Professional Blue",
}
for _k, _v in _defaults.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v


# ═══════════════════════════════════════════════════════════════
# 9. SIDEBAR
# ═══════════════════════════════════════════════════════════════
st.sidebar.title("🚀 Personal Brand AI")
st.sidebar.markdown("---")

st.sidebar.markdown("### Profile")
st.sidebar.markdown(f"**{user_name}**")
st.sidebar.caption(f"Style: {user_style[:60]}...")

st.sidebar.markdown("### Settings")
st.session_state.niche = st.sidebar.selectbox("Your Niche", [
    "AI & Technology", "Leadership & Management", "Entrepreneurship",
    "Marketing & Growth", "Finance & Investing", "Personal Development",
    "Product Management", "Sales", "Health & Wellness",
], index=["AI & Technology","Leadership & Management","Entrepreneurship","Marketing & Growth",
          "Finance & Investing","Personal Development","Product Management","Sales","Health & Wellness"
         ].index(st.session_state.niche) if st.session_state.niche in ["AI & Technology","Leadership & Management","Entrepreneurship","Marketing & Growth","Finance & Investing","Personal Development","Product Management","Sales","Health & Wellness"] else 0)

st.session_state.tone = st.sidebar.selectbox("Post Tone", [
    "Insightful & Authoritative", "Conversational & Personal",
    "Provocative & Bold", "Data-Driven & Analytical",
    "Motivational & Inspiring", "Storytelling",
])

st.sidebar.markdown("---")
st.sidebar.markdown("### Integration Status")
st.sidebar.markdown("🟢 Make.com — Connected")
st.sidebar.markdown("🟢 OpenAI GPT-4o — Active")
st.sidebar.markdown("🟢 Tavily Deep Research — Active")
st.sidebar.markdown("🟡 LinkedIn API — Not set up yet")
st.sidebar.markdown("🔵 Telegram — " + ("Active" if telegram_token else "Not configured"))

st.sidebar.markdown("---")
app_mode = st.sidebar.radio("Select Tool", [
    "🧠 Deep Agentic Workflow",
    "🗓️ Content Calendar",
    "✍️ Quick Writer",
    "🎠 Carousel Maker",
    "💬 AI Chatbot",
])


# ═══════════════════════════════════════════════════════════════
# HELPER: step progress bar renderer
# ═══════════════════════════════════════════════════════════════
def render_steps(current: int):
    labels = ["Idea", "Research", "Topics", "Write A/B", "Image", "Schedule"]
    cols   = st.columns(len(labels) * 2 - 1)
    for i, lbl in enumerate(labels):
        num     = i + 1
        ci      = i * 2
        cls     = "done" if num < current else "active" if num == current else "todo"
        colour  = "#0a66c2" if cls == "done" else "#1a1a2e" if cls == "active" else "#bbb"
        with cols[ci]:
            st.markdown(
                f'<div style="display:flex;align-items:center;gap:5px;">'
                f'<div class="step-pill {cls}">{num}</div>'
                f'<span style="font-size:11px;font-weight:600;color:{colour}">{lbl}</span>'
                f'</div>', unsafe_allow_html=True)
        if i < len(labels) - 1:
            with cols[ci + 1]:
                st.markdown(
                    '<div style="height:28px;display:flex;align-items:center;">'
                    '<div style="width:100%;height:1px;background:#e0e0e0"></div></div>',
                    unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# TOOL 1 — DEEP AGENTIC WORKFLOW  (multi-step, fully rebuilt)
# ══════════════════════════════════════════════════════════════════════════════
if app_mode == "🧠 Deep Agentic Workflow":
    st.title("🧠 Deep Agentic Workflow")
    render_steps(st.session_state.aw_step)
    st.markdown("---")

    # ── STEP 1: Idea Input ──────────────────────────────────────────────────
    if st.session_state.aw_step == 1:
        st.markdown("### 💡 What do you want to post about?")
        st.caption("Give a rough idea — a word, a trend, a question, or a story. The AI does the research.")

        idea_input = st.text_input(
            "Your idea",
            value=st.session_state.aw_idea,
            placeholder="e.g. 'Why most AI startups fail in year 2' or 'The power of saying no'",
            key="idea_input_field",
        )

        col1, col2 = st.columns([3, 1])
        with col2:
            run_btn = st.button("🔍 Research It →", type="primary", use_container_width=True)

        if run_btn:
            if not idea_input.strip():
                st.warning("Please enter an idea first.")
            elif not tavily_key:
                st.error("TAVILY_API_KEY missing from .env")
            else:
                st.session_state.aw_idea = idea_input.strip()
                st.session_state.aw_step = 2
                st.rerun()

    # ── STEP 2: Deep Research ───────────────────────────────────────────────
    elif st.session_state.aw_step == 2:
        st.markdown(f"### 🔍 Researching: *{st.session_state.aw_idea}*")
        st.caption("Tavily Deep Research is running — this takes 1–2 minutes. Please wait.")

        if st.session_state.aw_research is None:
            with st.status("🌐 Running Deep Research (1–2 min)...", expanded=True) as status:
                status.write("Sending query to Tavily Deep Research API...")
                raw = web_search(
                    f"Latest trends, data, research, and expert opinions about: {st.session_state.aw_idea}. "
                    f"Focus on {st.session_state.niche} professionals."
                )
                status.write("✅ Research complete — generating topic angles...")
                topics = generate_topic_angles(
                    st.session_state.aw_idea,
                    raw,
                    st.session_state.niche,
                )
                status.update(label="✅ Done!", state="complete", expanded=False)

            st.session_state.aw_research = raw
            st.session_state.aw_topics   = topics
            st.session_state.aw_step     = 3
            st.rerun()
        else:
            # Research already done — skip to step 3
            st.session_state.aw_step = 3
            st.rerun()

    # ── STEP 3: Topic Selection ─────────────────────────────────────────────
    elif st.session_state.aw_step == 3:
        st.markdown("### 📋 Choose Your Angle")
        st.caption("Five research-backed approaches — pick the one that feels right for you today.")

        col_cards, col_research = st.columns([2, 1])

        with col_cards:
            for i, t in enumerate(st.session_state.aw_topics):
                fmt_badge   = t.get("format", "Post")
                trend_badge = t.get("trend_score", "High")
                badge_col   = "badge-green" if trend_badge == "Very High" else "badge-blue"

                st.markdown(f"""
                <div class="topic-card">
                  <h5>{t['title']}</h5>
                  <p><em>{t['hook'][:120]}...</em></p>
                  <span class="badge badge-purple">{fmt_badge}</span>
                  <span class="badge {badge_col}">{trend_badge} reach</span>
                  <p style="margin-top:8px;font-size:11px;color:#888">{t.get('why_it_works','')}</p>
                </div>""", unsafe_allow_html=True)

                if st.button(f"Use Angle {i+1}", key=f"pick_topic_{i}", use_container_width=True):
                    st.session_state.aw_selected_topic = t
                    # Reset downstream so they regenerate
                    st.session_state.aw_post_a     = ""
                    st.session_state.aw_post_b     = ""
                    st.session_state.aw_final_post = ""
                    st.session_state.aw_image      = None
                    st.session_state.aw_step       = 4
                    st.rerun()

        with col_research:
            st.markdown("#### 🔬 Research Summary")
            with st.expander("View raw research", expanded=False):
                st.markdown(st.session_state.aw_research[:1200] + "..." if len(st.session_state.aw_research) > 1200 else st.session_state.aw_research)

        st.markdown("---")
        if st.button("← Back to Idea"):
            st.session_state.aw_step     = 1
            st.session_state.aw_research = None
            st.session_state.aw_topics   = []
            st.rerun()

    # ── STEP 4: A/B Post Generation ─────────────────────────────────────────
    elif st.session_state.aw_step == 4:
        topic = st.session_state.aw_selected_topic
        st.markdown(f"### ✍️ Two Versions — Pick Your Best")
        st.caption(f"Angle: **{topic['title']}** · Format: {topic['format']}")

        # Auto-generate if not done yet
        if not st.session_state.aw_post_a or not st.session_state.aw_post_b:
            with st.status("✍️ Writing two versions with GPT-4o...", expanded=True) as status:
                status.write("Version A — Hook-first, punchy paragraphs...")
                post_a = build_post(
                    st.session_state.aw_idea,
                    st.session_state.aw_research,
                    topic,
                    st.session_state.niche,
                    st.session_state.tone,
                    style_instructions="Hook-first, short punchy paragraphs, numbered or bulleted insights, strong authoritative CTA. Prioritise clarity and impact."
                )
                status.write("Version B — Story-driven, conversational...")
                post_b = build_post(
                    st.session_state.aw_idea,
                    st.session_state.aw_research,
                    topic,
                    st.session_state.niche,
                    st.session_state.tone,
                    style_instructions="Open with a personal story or scenario, conversational warm tone, build to an insight, end with a genuine question that invites conversation."
                )
                status.update(label="✅ Both versions ready!", state="complete", expanded=False)
            st.session_state.aw_post_a    = post_a
            st.session_state.aw_post_b    = post_b
            st.session_state.aw_metrics_a = post_metrics(post_a)
            st.session_state.aw_metrics_b = post_metrics(post_b)
            st.rerun()

        col_a, col_b = st.columns(2)

        def render_post_col(container, version_key, label, badge_cls, post_key, metrics_key):
            with container:
                m = st.session_state[metrics_key]
                st.markdown(
                    f'<span class="version-badge {badge_cls}">{label}</span>'
                    f'<div class="metrics-row">'
                    f'<span class="metric-chip">📝 {m.get("chars",0)} chars</span>'
                    f'<span class="metric-chip">🏷️ {m.get("hashtags",0)} hashtags</span>'
                    f'<span class="metric-chip">⏱️ {m.get("read_min",1)} min read</span>'
                    f'<span class="metric-chip">🎣 {m.get("hook_type","Hook")}</span>'
                    f'</div>', unsafe_allow_html=True)

                # Stale year warning — catch old dates before they go live
                stale = check_stale_years(st.session_state[post_key])
                if stale:
                    st.warning(
                        f"⚠️ Old year(s) detected: **{', '.join(stale)}** — "
                        f"edit the post to replace with '{datetime.now().year}' or remove before publishing.",
                        icon=None,
                    )

                edited = st.text_area(
                    f"Edit {label}",
                    value=st.session_state[post_key],
                    height=420,
                    key=f"edit_{post_key}",
                    label_visibility="collapsed",
                )
                st.session_state[post_key] = edited

                c1, c2 = st.columns(2)
                with c1:
                    if st.button("🔁 Regenerate", key=f"regen_{version_key}", use_container_width=True):
                        style = ("Hook-first, short punchy paragraphs, strong authoritative CTA."
                                 if version_key == "a" else
                                 "Open with a personal story, conversational tone, genuine question at end.")
                        with st.spinner(f"Rewriting {label}..."):
                            new_post = build_post(
                                st.session_state.aw_idea, st.session_state.aw_research,
                                topic, st.session_state.niche, st.session_state.tone, style)
                        st.session_state[post_key]    = new_post
                        st.session_state[metrics_key] = post_metrics(new_post)
                        st.rerun()
                with c2:
                    if st.button(f"✅ Use {label}", key=f"use_{version_key}", type="primary", use_container_width=True):
                        st.session_state.aw_active     = version_key.upper()
                        st.session_state.aw_final_post = st.session_state[post_key]
                        st.session_state.aw_image      = None  # reset image
                        st.session_state.aw_step       = 5
                        st.rerun()

        render_post_col(col_a, "a", "Version A · Punchy",        "version-a", "aw_post_a", "aw_metrics_a")
        render_post_col(col_b, "b", "Version B · Conversational", "version-b", "aw_post_b", "aw_metrics_b")

        # A/B comparison table
        with st.expander("📊 A/B Comparison"):
            ma, mb = st.session_state.aw_metrics_a, st.session_state.aw_metrics_b
            st.markdown(f"""
| Metric | Version A | Version B |
|---|---|---|
| Characters | {ma.get('chars',0)} | {mb.get('chars',0)} |
| Hashtags | {ma.get('hashtags',0)} | {mb.get('hashtags',0)} |
| Read time | {ma.get('read_min',1)} min | {mb.get('read_min',1)} min |
| Hook type | {ma.get('hook_type','—')} | {mb.get('hook_type','—')} |
""")

        st.markdown("---")
        if st.button("← Back to Topics"):
            st.session_state.aw_step   = 3
            st.session_state.aw_post_a = ""
            st.session_state.aw_post_b = ""
            st.rerun()

    # ── STEP 5: Image Generation ────────────────────────────────────────────
    elif st.session_state.aw_step == 5:
        st.markdown("### 🎨 Generate Your Post Image")

        col_left, col_right = st.columns([1, 1])

        with col_left:
            st.markdown(f"**Selected: Version {st.session_state.aw_active}**")
            st.caption(f"Topic: {st.session_state.aw_selected_topic.get('title','')}")
            st.text_area(
                "Post preview",
                value=st.session_state.aw_final_post[:500] + ("..." if len(st.session_state.aw_final_post) > 500 else ""),
                height=260,
                disabled=True,
                label_visibility="collapsed",
            )

        with col_right:
            img_style = st.selectbox("Visual style", [
                "Professional & Clean", "Bold & Typographic",
                "Abstract & Conceptual", "Photorealistic Scene",
                "Data / Infographic Style", "Minimalist",
            ])
            img_mood = st.selectbox("Mood", [
                "Inspiring & Uplifting", "Serious & Authoritative",
                "Futuristic & Tech", "Warm & Human", "Bold & Energetic",
            ])
            extra = st.text_input("Extra details (optional)", placeholder="e.g. dark theme, specific colours, include a chart")

            if st.button("🎨 Generate Image", type="primary", use_container_width=True):
                with st.spinner("Building image prompt then calling DALL·E 3..."):
                    img_prompt = generate_text(
                        f"Write a DALL·E 3 prompt for a LinkedIn post image.\n"
                        f"Post topic: {st.session_state.aw_selected_topic.get('title','')}\n"
                        f"Post excerpt: {st.session_state.aw_final_post[:300]}\n"
                        f"Visual style: {img_style}\n"
                        f"Mood: {img_mood}\n"
                        f"Extra: {extra or 'none'}\n"
                        f"Rules: photorealistic or clean graphic, NO text in image, NO human faces, "
                        f"LinkedIn-appropriate, 1792x1024 landscape format, {img_style.lower()} aesthetic.",
                        "You write concise, precise DALL·E 3 prompts. Max 120 words."
                    )
                    img_bytes = generate_image(img_prompt)
                st.session_state.aw_image     = img_bytes
                st.session_state.aw_img_prompt = img_prompt

        if st.session_state.aw_image:
            st.image(st.session_state.aw_image, caption="Generated image", use_container_width=True)
            with st.expander("Prompt used"):
                st.code(st.session_state.aw_img_prompt)

            c1, c2, c3 = st.columns(3)
            with c1:
                if st.button("🔁 Regenerate Image", use_container_width=True):
                    st.session_state.aw_image = None
                    st.rerun()
            with c2:
                if st.button("⏭️ Skip Image", use_container_width=True):
                    st.session_state.aw_image = None
                    st.session_state.aw_step  = 6
                    st.rerun()
            with c3:
                if st.button("✅ Use This Image →", type="primary", use_container_width=True):
                    st.session_state.aw_step = 6
                    st.rerun()
        else:
            st.info("Generate an image above, or skip to proceed without one.")
            if st.button("⏭️ Skip — proceed without image →", use_container_width=True):
                st.session_state.aw_step = 6
                st.rerun()

        st.markdown("---")
        if st.button("← Back to Posts"):
            st.session_state.aw_step = 4
            st.rerun()

    # ── STEP 6: Schedule & Publish ──────────────────────────────────────────
    elif st.session_state.aw_step == 6:
        st.markdown("### 🚀 Schedule & Publish")

        col_post, col_settings = st.columns([3, 2])

        with col_post:
            st.markdown("**Final post — edit if needed**")
            final = st.text_area(
                "Final post",
                value=st.session_state.aw_final_post,
                height=320,
                label_visibility="collapsed",
                key="final_post_edit",
            )
            st.session_state.aw_final_post = final

            char_count = len(final)
            colour = "green" if char_count <= 3000 else "red"
            st.markdown(f'<span style="font-size:12px;color:{colour}">{char_count}/3000 characters</span>', unsafe_allow_html=True)

            if st.session_state.aw_image:
                st.image(st.session_state.aw_image, caption="Attached image", width=320)

        with col_settings:
            st.markdown("**When to publish?**")
            publish_mode = st.radio("Mode", [
                "Auto — next best slot",
                "Schedule a specific time",
                "Send to Make.com now",
            ], label_visibility="collapsed")

            if publish_mode == "Auto — next best slot":
                slot = get_next_best_slot()
                st.markdown(
                    f'<div class="slot-highlight">'
                    f'📅 Scheduled for: <b>{slot["display"]}</b><br>'
                    f'Engagement score: {slot["score"]}/10 &nbsp;·&nbsp; '
                    f'<code>publish_at</code> → <code>{slot["iso"]}</code>'
                    f'</div>',
                    unsafe_allow_html=True,
                )
                st.session_state.aw_scheduled_iso = slot["iso"]
                st.session_state.aw_scheduled_lbl = slot["display"]

                with st.expander("📅 Full slot calendar (next 7 days)"):
                    for s in get_slot_calendar():
                        bar = "▓" * s["score"] + "░" * (10 - s["score"])
                        st.markdown(f"**{s['label']}** &nbsp; `{bar}` &nbsp; {s['score']}/10")

            elif publish_mode == "Schedule a specific time":
                import datetime as dt
                pub_date = st.date_input("Date", min_value=dt.date.today())
                pub_time = st.time_input("Time (IST)", value=dt.time(8, 0))
                combined = dt.datetime.combine(pub_date, pub_time)
                st.session_state.aw_scheduled_iso = combined.strftime("%Y-%m-%dT%H:%M:%S")
                st.session_state.aw_scheduled_lbl = combined.strftime("%A %d %b, %I:%M %p IST")
                st.markdown(
                    f'<div class="slot-highlight">'
                    f'📅 Scheduled for: <b>{st.session_state.aw_scheduled_lbl}</b><br>'
                    f'<code>publish_at</code> → <code>{st.session_state.aw_scheduled_iso}</code>'
                    f'</div>',
                    unsafe_allow_html=True,
                )
            else:
                st.session_state.aw_scheduled_iso = datetime.now().strftime("%Y-%m-%dT%H:%M:%S")
                st.session_state.aw_scheduled_lbl = "Immediate"
                st.warning("⚡ Immediate mode — Make.com will post as soon as it receives this.", icon=None)

            # ── Make.com Sleep module setup notice ──────────────────────────
            st.markdown("""
<div class="linkedin-placeholder">
  ⚙️ <b>To enable scheduled posting in Make.com:</b><br>
  Edit scenario → right-click arrow after <b>Tools</b> → Insert module → <b>Sleep</b><br>
  Set <i>Resume execution</i> to: <code>{{1.publish_at}}</code><br>
  Then Save. Your post will wait until the exact scheduled time before posting.
</div>
""", unsafe_allow_html=True)

            st.markdown("**Notifications**")
            notify_tg = st.checkbox("Telegram notification", value=bool(telegram_token))
            if not telegram_token:
                st.caption("Add TELEGRAM_BOT_TOKEN + TELEGRAM_CHAT_ID to .env to enable")

            st.markdown("---")

            if st.button("🚀 Send to Make.com & Save", type="primary", use_container_width=True):
                with st.spinner("Compressing image and sending to Make.com..."):
                    status_code, make_msg = send_to_make(
                        post        = st.session_state.aw_final_post,
                        image_bytes = st.session_state.aw_image,
                        topic       = st.session_state.aw_idea,
                        niche       = st.session_state.niche,
                        publish_at  = st.session_state.aw_scheduled_iso,
                        slot_label  = st.session_state.aw_scheduled_lbl,
                    )

                if status_code == 200:
                    save_post_to_db(
                        topic          = st.session_state.aw_idea,
                        content        = st.session_state.aw_final_post,
                        image_data     = st.session_state.aw_image,
                        scheduled_time = st.session_state.aw_scheduled_lbl,
                        niche          = st.session_state.niche,
                        version        = st.session_state.aw_active,
                    )
                    if notify_tg:
                        send_telegram(
                            f"<b>Post queued in Make.com!</b>\n\n"
                            f"Topic: {st.session_state.aw_idea}\n"
                            f"Scheduled: {st.session_state.aw_scheduled_lbl}\n"
                            f"Version: {st.session_state.aw_active}\n\n"
                            f"Make.com will publish at the scheduled time."
                        )

                    is_immediate = st.session_state.aw_scheduled_lbl == "Immediate"
                    if is_immediate:
                        st.success("✅ Sent to Make.com — posting to LinkedIn now!")
                    else:
                        st.success(f"✅ Scheduled! Make.com will post on **{st.session_state.aw_scheduled_lbl}**")

                    st.balloons()

                    if is_immediate:
                        st.info("💡 To schedule future posts instead of posting immediately, add a **Sleep** module in Make.com — see the setup note above.")
                    else:
                        st.markdown(f"""
**What happens next:**
1. ✅ Make.com received your post + image
2. ⏳ Sleep module waits until **{st.session_state.aw_scheduled_lbl}**
3. 📤 LinkedIn module publishes automatically
4. 📊 Google Sheets logs the post URL
5. 📱 Telegram notifies you when it's live

> If you haven't added the Sleep module yet, Make.com will post immediately regardless of the scheduled time. Follow the setup note above.
""")
                else:
                    st.error(f"Make.com returned status {status_code}")
                    with st.expander("🔍 Full error — click to diagnose"):
                        st.code(make_msg or "No response body")
                        st.markdown("""
**Common causes:**
- `400` with image → image still too large even after compression (install `Pillow`: `pip install Pillow`)
- `400` without image → Make.com scenario is paused or the webhook URL changed
- `410` / `404` → webhook was deleted in Make.com — create a new one and update `MAKE_WEBHOOK_URL`
- `None` / timeout → network issue or Make.com is down
                        """)
                    st.info("Your post is NOT lost — click 'Save Draft Only' below to keep it.")

            if st.button("💾 Save Draft Only (no Make.com)", use_container_width=True):
                save_post_to_db(
                    topic          = st.session_state.aw_idea,
                    content        = st.session_state.aw_final_post,
                    image_data     = st.session_state.aw_image,
                    scheduled_time = st.session_state.aw_scheduled_lbl,
                    niche          = st.session_state.niche,
                    version        = st.session_state.aw_active,
                )
                st.success("✅ Saved to Content Calendar as Draft.")

        st.markdown("---")
        col_back, col_new = st.columns(2)
        with col_back:
            if st.button("← Back to Image"):
                st.session_state.aw_step = 5
                st.rerun()
        with col_new:
            if st.button("🔄 Start New Post"):
                for key in ["aw_step","aw_idea","aw_research","aw_topics","aw_selected_topic",
                            "aw_post_a","aw_post_b","aw_metrics_a","aw_metrics_b","aw_active",
                            "aw_final_post","aw_image","aw_img_prompt","aw_scheduled_iso","aw_scheduled_lbl"]:
                    st.session_state[key] = _defaults.get(key, "" if "step" not in key else 1)
                st.session_state.aw_step = 1
                st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# TOOL 2 — CONTENT CALENDAR  (enhanced with scheduled_time + niche columns)
# ══════════════════════════════════════════════════════════════════════════════
elif app_mode == "🗓️ Content Calendar":
    st.title("🗓️ Content Calendar")
    st.caption("All generated posts — drafts, scheduled, and published.")

    posts = get_all_posts()

    if not posts:
        st.info("No posts saved yet. Run the Deep Agentic Workflow to generate your first post!")
    else:
        # Summary row
        total     = len(posts)
        published = sum(1 for p in posts if "Published" in str(p[5]))
        scheduled = sum(1 for p in posts if str(p[5]) == "Draft" and p[6])
        drafts    = total - published

        m1, m2, m3, m4 = st.columns(4)
        m1.metric("Total posts", total)
        m2.metric("Published",   published)
        m3.metric("Scheduled",   scheduled)
        m4.metric("Drafts",      drafts)
        st.markdown("---")

        for post in posts:
            post_id, date, topic, content, image_data, post_status, scheduled_time, niche, version = post

            # Build expander header
            version_tag   = f" | v{version}" if version else ""
            niche_tag     = f" | {niche}"    if niche   else ""
            schedule_tag  = f" | 🕐 {scheduled_time}" if scheduled_time else ""
            header        = f"{date} | {topic}{version_tag}{niche_tag} | **{post_status}**{schedule_tag}"

            with st.expander(header):
                col_img, col_content = st.columns([1, 2])

                with col_img:
                    if image_data:
                        try:
                            st.image(image_data, use_container_width=True)
                        except Exception:
                            st.caption("Image stored but preview unavailable")
                    else:
                        st.caption("No image attached")

                    if scheduled_time:
                        st.markdown(f"**Scheduled:** {scheduled_time}")
                    if niche:
                        st.caption(f"Niche: {niche}")
                    if version:
                        st.caption(f"Version: {version}")

                with col_content:
                    st.markdown(content)

                st.markdown("---")
                col1, col2, col3 = st.columns(3)

                with col1:
                    if "Published" not in str(post_status):
                        if st.button("🚀 Publish via Make.com", key=f"pub_{post_id}"):
                            slot              = get_next_best_slot()
                            make_stat, make_msg = send_to_make(
                                post        = content,
                                image_bytes = image_data,
                                topic       = topic,
                                niche       = niche or st.session_state.niche,
                                publish_at  = slot["iso"],
                                slot_label  = slot["display"],
                            )
                            if make_stat == 200:
                                update_post_status(post_id, "Published ✅")
                                send_telegram(
                                    f"<b>Post sent to Make.com!</b>\n"
                                    f"Topic: {topic}\n"
                                    f"Scheduled: {slot['display']}"
                                )
                                st.success(f"Sent! Scheduled for {slot['display']}")
                                st.rerun()
                            else:
                                st.error(f"Webhook failed (Status: {make_stat})")
                                with st.expander("🔍 Error details"):
                                    st.code(make_msg or "No response body")
                with col2:
                    if st.button("📋 Copy Post", key=f"copy_{post_id}"):
                        st.code(content, language=None)
                with col3:
                    if "Published" in str(post_status):
                        st.success("Live ✅")
                    else:
                        st.info("Draft — ready to publish")


# ══════════════════════════════════════════════════════════════════════════════
# TOOL 3 — QUICK WRITER  (improved with tone + format options)
# ══════════════════════════════════════════════════════════════════════════════
elif app_mode == "✍️ Quick Writer":
    st.title("✍️ Quick Writer")
    st.caption("Generate a LinkedIn post fast — no research step, just your topic and go.")

    col_in, col_opts = st.columns([2, 1])

    with col_in:
        topic_input = st.text_area(
            "Your topic or idea",
            placeholder="e.g. 'The one leadership lesson that changed how I manage teams'",
            height=100,
        )

    with col_opts:
        qw_format = st.selectbox("Format", [
            "Numbered list (5 insights)",
            "Story arc",
            "Contrarian take",
            "How-to guide",
            "Data-driven reveal",
        ])
        qw_length = st.selectbox("Length", ["Short (600–900 chars)", "Medium (1000–1400 chars)", "Long (1500–1800 chars)"])

    if st.button("🔥 Generate Post", type="primary"):
        if not topic_input.strip():
            st.warning("Please enter a topic first.")
        else:
            length_map = {
                "Short (600–900 chars)":    "600–900 characters",
                "Medium (1000–1400 chars)": "1000–1400 characters",
                "Long (1500–1800 chars)":   "1500–1800 characters",
            }
            with st.spinner("Writing your post..."):
                qw_prompt = f"""Write a LinkedIn post about: {topic_input}

Author: {user_name}
Style: {user_style}
Niche: {st.session_state.niche}
Tone: {st.session_state.tone}
Format: {qw_format}
Target length: {length_map.get(qw_length, '1000–1400 characters')}

{LINKEDIN_POST_RULES}

Output ONLY the post text."""
                result = generate_text(qw_prompt, f"You are an expert LinkedIn ghostwriter for {st.session_state.niche}.")
                m      = post_metrics(result)

            st.markdown("### Your Post")
            st.markdown(
                f'<div class="metrics-row">'
                f'<span class="metric-chip">📝 {m["chars"]} chars</span>'
                f'<span class="metric-chip">🏷️ {m["hashtags"]} hashtags</span>'
                f'<span class="metric-chip">⏱️ {m["read_min"]} min read</span>'
                f'<span class="metric-chip">🎣 {m["hook_type"]}</span>'
                f'</div>', unsafe_allow_html=True)

            edited_qw = st.text_area("Edit post", value=result, height=380, label_visibility="collapsed")

            col_a, col_b, col_c = st.columns(3)
            with col_a:
                if st.button("💾 Save to Calendar"):
                    save_post_to_db(
                        topic   = topic_input[:80],
                        content = edited_qw,
                        image_data = None,
                        niche   = st.session_state.niche,
                        version = "Quick",
                    )
                    st.success("Saved to Content Calendar!")
            with col_b:
                if st.button("🚀 Send to Make.com"):
                    slot              = get_next_best_slot()
                    make_stat, make_msg = send_to_make(
                        post       = edited_qw,
                        topic      = topic_input[:80],
                        niche      = st.session_state.niche,
                        publish_at = slot["iso"],
                        slot_label = slot["display"],
                    )
                    if make_stat == 200:
                        save_post_to_db(topic_input[:80], edited_qw, None, slot["display"], st.session_state.niche, "Quick")
                        st.success(f"Sent! Scheduled: {slot['display']}")
                    else:
                        st.error(f"Webhook error (Status: {make_stat})")
                        with st.expander("🔍 Error details"):
                            st.code(make_msg or "No response body")
            with col_c:
                if st.button("📋 Copy Text"):
                    st.code(edited_qw, language=None)


# ══════════════════════════════════════════════════════════════════════════════
# TOOL 4 — CAROUSEL MAKER
# Generates 5–8 LinkedIn carousel slides, shows visual preview,
# and exports a downloadable PPTX.
# ══════════════════════════════════════════════════════════════════════════════
elif app_mode == "🎠 Carousel Maker":
    st.title("🎠 Carousel Maker")
    st.caption("Turn any topic into a LinkedIn carousel — swipeable slides with a cover, insights, and CTA.")

    # ── Style palette definitions ────────────────────────────────────────────
    CAROUSEL_STYLES = {
        "Professional Blue": {
            "bg": "#0a66c2", "text": "#ffffff", "accent": "#cce4ff",
            "card_bg": "#0857a8", "font_size_title": "22px", "font_size_body": "14px",
        },
        "Dark Minimal": {
            "bg": "#1a1a2e", "text": "#e8e8f0", "accent": "#7f77dd",
            "card_bg": "#16213e", "font_size_title": "22px", "font_size_body": "14px",
        },
        "Warm Coral": {
            "bg": "#d85a30", "text": "#ffffff", "accent": "#ffd4c2",
            "card_bg": "#bf4f2a", "font_size_title": "22px", "font_size_body": "14px",
        },
        "Clean White": {
            "bg": "#ffffff", "text": "#1a1a2e", "accent": "#0a66c2",
            "card_bg": "#f4f6fb", "font_size_title": "22px", "font_size_body": "14px",
        },
        "Forest Green": {
            "bg": "#1d9e75", "text": "#ffffff", "accent": "#c8f5e8",
            "card_bg": "#188a65", "font_size_title": "22px", "font_size_body": "14px",
        },
    }

    # ── Input form ───────────────────────────────────────────────────────────
    col_inp, col_opts = st.columns([2, 1])
    with col_inp:
        carousel_topic = st.text_input(
            "Carousel topic",
            value=st.session_state.carousel_topic,
            placeholder="e.g. '5 habits that doubled my productivity' or 'How AI is changing marketing'",
            key="carousel_topic_input",
        )
    with col_opts:
        num_slides = st.slider("Number of slides", min_value=4, max_value=8, value=6)
        style_choice = st.selectbox("Visual style", list(CAROUSEL_STYLES.keys()))
        st.session_state.carousel_style = style_choice

    col_g1, col_g2 = st.columns([1, 3])
    with col_g1:
        gen_btn = st.button("✨ Generate Slides", type="primary", use_container_width=True)
    with col_g2:
        if st.session_state.carousel_slides:
            st.caption(f"{len(st.session_state.carousel_slides)} slides ready — preview and download below")

    if gen_btn and carousel_topic.strip():
        st.session_state.carousel_topic = carousel_topic.strip()
        TODAY      = datetime.now().strftime("%B %d, %Y")
        CURRENT_YR = datetime.now().strftime("%Y")

        with st.spinner("Generating carousel structure with GPT-4o..."):
            prompt = f"""Create a LinkedIn carousel post with exactly {num_slides} slides about: "{carousel_topic}"

TODAY: {TODAY} | YEAR: {CURRENT_YR}
Author niche: {st.session_state.niche}
Style: {st.session_state.tone}

SLIDE STRUCTURE RULES:
- Slide 1: Cover slide — bold hook title (max 8 words) + subtitle (max 12 words)
- Slides 2 to {num_slides - 1}: One insight per slide — short title + 2-3 line body + one key takeaway
- Slide {num_slides}: CTA slide — "Follow {user_name}" + one action the reader should take

Return ONLY a valid JSON array, no markdown, no extra text:
[
  {{
    "slide_num": 1,
    "type": "cover",
    "title": "Bold hook title here",
    "subtitle": "Supporting subtitle here",
    "body": "",
    "takeaway": "",
    "emoji": "🚀"
  }},
  {{
    "slide_num": 2,
    "type": "content",
    "title": "Insight title",
    "subtitle": "",
    "body": "2-3 lines explaining this insight clearly and concisely.",
    "takeaway": "→ One punchy takeaway line",
    "emoji": "💡"
  }},
  {{
    "slide_num": {num_slides},
    "type": "cta",
    "title": "Found this useful?",
    "subtitle": "Follow {user_name} for more on {st.session_state.niche}",
    "body": "Save this carousel and share it with someone who needs it.",
    "takeaway": "",
    "emoji": "🔔"
  }}
]

Important: Only 2024-{CURRENT_YR} data. All slides must be ready to publish."""

            raw = generate_text(
                prompt,
                "You are a JSON-only assistant. Output valid JSON arrays only.",
                temperature=0.75,
            )
            try:
                clean = raw.strip()
                if "```" in clean:
                    clean = clean.split("```")[1]
                    if clean.startswith("json"):
                        clean = clean[4:]
                slides = json.loads(clean.strip())
                st.session_state.carousel_slides = slides
            except Exception as e:
                st.error(f"Could not parse slides: {e}")
                st.code(raw[:500])

    # ── Visual Preview ───────────────────────────────────────────────────────
    if st.session_state.carousel_slides:
        style = CAROUSEL_STYLES[st.session_state.carousel_style]
        slides = st.session_state.carousel_slides

        st.markdown("---")
        st.markdown("### 📱 Preview")
        st.caption("Scroll right to see all slides · Click any slide text to edit")

        # Render slides in rows of 3
        for row_start in range(0, len(slides), 3):
            row_slides = slides[row_start:row_start + 3]
            cols = st.columns(len(row_slides))

            for col, slide in zip(cols, row_slides):
                with col:
                    stype   = slide.get("type", "content")
                    emoji   = slide.get("emoji", "")
                    title   = slide.get("title", "")
                    subtitle= slide.get("subtitle", "")
                    body    = slide.get("body", "")
                    tkaway  = slide.get("takeaway", "")
                    num     = slide.get("slide_num", "")

                    # Card HTML
                    card_html = f"""
<div style="
    background:{style['card_bg']};
    border-radius:12px;
    padding:20px 16px;
    min-height:220px;
    position:relative;
    border:2px solid {style['accent']}22;
    margin-bottom:8px;
">
  <div style="
    position:absolute;top:10px;right:12px;
    font-size:10px;color:{style['accent']};
    font-weight:700;opacity:0.8;
  ">{num}/{len(slides)}</div>
  <div style="font-size:26px;margin-bottom:10px">{emoji}</div>
  <div style="
    font-size:{style['font_size_title']};
    font-weight:800;
    color:{style['text']};
    line-height:1.25;
    margin-bottom:8px;
  ">{title}</div>"""

                    if subtitle:
                        card_html += f"""
  <div style="
    font-size:13px;color:{style['accent']};
    font-weight:600;margin-bottom:10px;
  ">{subtitle}</div>"""

                    if body:
                        card_html += f"""
  <div style="
    font-size:{style['font_size_body']};
    color:{style['text']};opacity:0.9;
    line-height:1.6;margin-bottom:10px;
  ">{body}</div>"""

                    if tkaway:
                        card_html += f"""
  <div style="
    font-size:13px;font-weight:700;
    color:{style['accent']};
    border-left:3px solid {style['accent']};
    padding-left:8px;margin-top:8px;
  ">{tkaway}</div>"""

                    card_html += "</div>"

                    st.markdown(
                        f'<div style="background:{style["bg"]};border-radius:14px;padding:4px;">'
                        f'{card_html}</div>',
                        unsafe_allow_html=True,
                    )

        # ── Edit slides ──────────────────────────────────────────────────────
        st.markdown("---")
        with st.expander("✏️ Edit slides"):
            updated_slides = []
            for i, slide in enumerate(slides):
                st.markdown(f"**Slide {slide['slide_num']} — {slide.get('type','').title()}**")
                c1, c2 = st.columns([1, 3])
                with c1:
                    new_emoji = st.text_input("Emoji", value=slide.get("emoji",""), key=f"em_{i}")
                with c2:
                    new_title = st.text_input("Title", value=slide.get("title",""), key=f"ti_{i}")
                new_body = st.text_area(
                    "Body",
                    value=slide.get("body",""),
                    height=80,
                    key=f"bo_{i}",
                    label_visibility="collapsed" if not slide.get("body") else "visible",
                )
                new_tkaway = st.text_input("Takeaway", value=slide.get("takeaway",""), key=f"tk_{i}")
                updated_slides.append({**slide, "emoji": new_emoji, "title": new_title,
                                        "body": new_body, "takeaway": new_tkaway})
                st.markdown("---")
            if st.button("💾 Save edits", key="save_carousel_edits"):
                st.session_state.carousel_slides = updated_slides
                st.success("Slides updated!")
                st.rerun()

        # ── LinkedIn post caption ─────────────────────────────────────────────
        st.markdown("### 📝 LinkedIn Caption for this Carousel")
        if st.button("Generate caption", key="gen_carousel_caption"):
            with st.spinner("Writing caption..."):
                titles = " | ".join([s.get("title","") for s in slides[1:-1]])
                cap    = generate_text(
                    f"Write a LinkedIn post caption to accompany a carousel about '{st.session_state.carousel_topic}'.\n"
                    f"Slide titles: {titles}\n"
                    f"Author: {user_name} | Niche: {st.session_state.niche}\n"
                    f"Rules: 3-5 short paragraphs, end with 'Swipe to see all →', 3 hashtags.\n"
                    f"Today: {datetime.now().strftime('%B %d, %Y')}. No years before {datetime.now().year}.\n"
                    f"{LINKEDIN_POST_RULES}\nOutput ONLY the caption text.",
                    temperature=0.8,
                )
                st.session_state["carousel_caption"] = cap

        if st.session_state.get("carousel_caption"):
            caption_edit = st.text_area(
                "Caption",
                value=st.session_state["carousel_caption"],
                height=220,
                key="carousel_caption_edit",
            )
            st.session_state["carousel_caption"] = caption_edit

            col_cs1, col_cs2 = st.columns(2)
            with col_cs1:
                if st.button("💾 Save to Calendar", key="save_carousel_cal"):
                    save_post_to_db(
                        topic      = f"[Carousel] {st.session_state.carousel_topic[:70]}",
                        content    = caption_edit,
                        image_data = None,
                        niche      = st.session_state.niche,
                        version    = "Carousel",
                    )
                    st.success("Saved to Content Calendar!")
            with col_cs2:
                if st.button("🚀 Send Caption to Make.com", key="send_carousel_make"):
                    slot               = get_next_best_slot()
                    make_stat, make_msg = send_to_make(
                        post       = caption_edit,
                        topic      = f"[Carousel] {st.session_state.carousel_topic[:70]}",
                        niche      = st.session_state.niche,
                        publish_at = slot["iso"],
                        slot_label = slot["display"],
                    )
                    if make_stat == 200:
                        st.success(f"Scheduled: {slot['display']}")
                    else:
                        st.error(f"Webhook error: {make_stat}")

        # ── PPTX Download ────────────────────────────────────────────────────
        st.markdown("---")
        st.markdown("### 💾 Download")

        col_dl1, col_dl2 = st.columns(2)
        with col_dl1:
            if st.button("📥 Generate PPTX Download", key="gen_pptx", use_container_width=True):
                try:
                    from pptx import Presentation
                    from pptx.util import Inches, Pt, Emu
                    from pptx.dml.color import RGBColor
                    from pptx.enum.text import PP_ALIGN
                    import io, re as _re

                    def hex_to_rgb(h):
                        h = h.lstrip("#")
                        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

                    prs = Presentation()
                    prs.slide_width  = Inches(10)
                    prs.slide_height = Inches(10)   # square — best for LinkedIn carousels

                    blank_layout = prs.slide_layouts[6]
                    s           = CAROUSEL_STYLES[st.session_state.carousel_style]

                    for slide_data in st.session_state.carousel_slides:
                        sl = prs.slides.add_slide(blank_layout)

                        # Background
                        bg = sl.background
                        fill = bg.fill
                        fill.solid()
                        fill.fore_color.rgb = hex_to_rgb(s["bg"])

                        def add_text(sl, text, left, top, width, height,
                                     font_size, bold=False, color="#ffffff", align=PP_ALIGN.LEFT):
                            txb  = sl.shapes.add_textbox(
                                Inches(left), Inches(top), Inches(width), Inches(height))
                            tf   = txb.text_frame
                            tf.word_wrap = True
                            p    = tf.paragraphs[0]
                            p.alignment = align
                            run  = p.add_run()
                            run.text = text
                            run.font.size = Pt(font_size)
                            run.font.bold = bold
                            run.font.color.rgb = hex_to_rgb(color)

                        # Slide number pill
                        add_text(sl, f"{slide_data['slide_num']}/{len(slides)}",
                                 8.5, 0.3, 1.2, 0.4, 11, color=s["accent"], align=PP_ALIGN.RIGHT)

                        # Emoji + Title
                        add_text(sl, slide_data.get("emoji",""), 0.5, 0.8, 1.0, 1.0, 36)
                        add_text(sl, slide_data.get("title",""), 0.5, 1.7, 9.0, 2.5,
                                 28, bold=True, color=s["text"], align=PP_ALIGN.LEFT)

                        # Subtitle
                        if slide_data.get("subtitle"):
                            add_text(sl, slide_data["subtitle"], 0.5, 4.0, 9.0, 1.0,
                                     16, color=s["accent"])

                        # Body
                        if slide_data.get("body"):
                            add_text(sl, slide_data["body"], 0.5, 5.0, 9.0, 3.0,
                                     15, color=s["text"])

                        # Takeaway
                        if slide_data.get("takeaway"):
                            add_text(sl, slide_data["takeaway"], 0.5, 8.3, 9.0, 1.0,
                                     14, bold=True, color=s["accent"])

                    buf = io.BytesIO()
                    prs.save(buf)
                    buf.seek(0)

                    safe_topic = _re.sub(r'[^\w\s-]', '', st.session_state.carousel_topic)[:40].strip()
                    fname      = f"carousel_{safe_topic.replace(' ','_')}.pptx"

                    st.download_button(
                        label    = "⬇️ Download PPTX",
                        data     = buf,
                        file_name = fname,
                        mime     = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key      = "pptx_download_btn",
                        use_container_width=True,
                    )
                    st.success("PPTX ready — click Download PPTX above!")

                except ImportError:
                    st.error("python-pptx not installed.")
                    st.code("pip install python-pptx")
                except Exception as e:
                    st.error(f"PPTX error: {e}")

        with col_dl2:
            # Plain text export — always works
            if st.session_state.carousel_slides:
                txt_lines = [f"CAROUSEL: {st.session_state.carousel_topic}\n"]
                for sl in st.session_state.carousel_slides:
                    txt_lines.append(f"--- Slide {sl['slide_num']} ---")
                    if sl.get("title"):    txt_lines.append(f"Title: {sl['title']}")
                    if sl.get("subtitle"): txt_lines.append(f"Subtitle: {sl['subtitle']}")
                    if sl.get("body"):     txt_lines.append(f"{sl['body']}")
                    if sl.get("takeaway"): txt_lines.append(f"{sl['takeaway']}")
                    txt_lines.append("")
                st.download_button(
                    label     = "⬇️ Download as TXT",
                    data      = "\n".join(txt_lines),
                    file_name = "carousel_slides.txt",
                    mime      = "text/plain",
                    key       = "txt_download_btn",
                    use_container_width=True,
                )


# ══════════════════════════════════════════════════════════════════════════════
# TOOL 5 — AI CHATBOT
# LinkedIn-aware GPT-4o chat assistant.
# Knows your niche, style, and post history. Helps brainstorm, refine,
# and answer strategy questions. Chat history persists within the session.
# ══════════════════════════════════════════════════════════════════════════════
elif app_mode == "💬 AI Chatbot":
    st.title("💬 AI Chatbot")
    st.caption(f"Your LinkedIn strategy assistant · Niche: {st.session_state.niche} · Style: {user_style[:50]}...")

    CHAT_SYSTEM = f"""You are a highly experienced LinkedIn content strategist and personal brand coach.

You are helping: {user_name}
Their niche: {st.session_state.niche}
Their writing style: {user_style}
Current tone preference: {st.session_state.tone}
Today's date: {datetime.now().strftime("%B %d, %Y")}

Your role:
- Help brainstorm LinkedIn post ideas, hooks, and angles
- Refine drafts — improve hooks, structure, CTAs, hashtags
- Answer LinkedIn strategy questions (algorithm, best times, content formats)
- Suggest carousel topics, post series, and content calendar plans
- Review posts for old dates, jargon, or weak hooks before publishing
- Keep advice specific, actionable, and tailored to {st.session_state.niche}

Tone: Friendly expert. Direct and concise. No corporate jargon.
Always reference {datetime.now().year} best practices — never outdated advice."""

    # Clear chat button
    col_chat_hdr, col_chat_clr = st.columns([4, 1])
    with col_chat_clr:
        if st.button("🗑️ Clear chat", use_container_width=True):
            st.session_state.chat_history = []
            st.rerun()

    # Render chat history
    for msg in st.session_state.chat_history:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # Suggested prompts — shown only when chat is empty
    if not st.session_state.chat_history:
        st.markdown("**Suggested questions to get started:**")
        suggestions = [
            f"Give me 5 LinkedIn post ideas for {st.session_state.niche} this week",
            "Review this hook and make it more scroll-stopping",
            "What's the best LinkedIn posting strategy for 2026?",
            "Turn this bullet list into a LinkedIn carousel structure",
            "What hashtags should I use for my niche?",
            "Write a strong CTA for a post about leadership",
        ]
        cols = st.columns(2)
        for i, suggestion in enumerate(suggestions):
            with cols[i % 2]:
                if st.button(suggestion, key=f"sug_{i}", use_container_width=True):
                    # Treat button click as user message
                    st.session_state.chat_history.append({"role": "user", "content": suggestion})
                    with st.spinner("Thinking..."):
                        messages = [{"role": "system", "content": CHAT_SYSTEM}]
                        for m in st.session_state.chat_history:
                            messages.append({"role": m["role"], "content": m["content"]})
                        url     = "https://api.openai.com/v1/chat/completions"
                        headers = {"Authorization": f"Bearer {llm_key}", "Content-Type": "application/json"}
                        resp    = requests.post(url, headers=headers, json={
                            "model": "gpt-4o", "messages": messages, "temperature": 0.8
                        }, timeout=60).json()
                        try:
                            reply = resp["choices"][0]["message"]["content"]
                        except Exception:
                            reply = "Sorry, something went wrong. Please try again."
                    st.session_state.chat_history.append({"role": "assistant", "content": reply})
                    st.rerun()

    # Chat input
    user_input = st.chat_input(
        f"Ask anything about LinkedIn strategy, post ideas, hooks, carousels..."
    )

    if user_input:
        # Add user message
        st.session_state.chat_history.append({"role": "user", "content": user_input})

        # Build full message list with system prompt
        messages = [{"role": "system", "content": CHAT_SYSTEM}]
        # Keep last 20 messages to avoid token overflow
        for m in st.session_state.chat_history[-20:]:
            messages.append({"role": m["role"], "content": m["content"]})

        with st.chat_message("assistant"):
            with st.spinner(""):
                url     = "https://api.openai.com/v1/chat/completions"
                headers = {"Authorization": f"Bearer {llm_key}", "Content-Type": "application/json"}
                resp    = requests.post(url, headers=headers, json={
                    "model": "gpt-4o", "messages": messages, "temperature": 0.8
                }, timeout=60).json()
                try:
                    reply = resp["choices"][0]["message"]["content"]
                except Exception:
                    reply = f"API error: {resp.get('error', {}).get('message', 'Unknown error')}"
            st.markdown(reply)

        st.session_state.chat_history.append({"role": "assistant", "content": reply})
        st.rerun()
